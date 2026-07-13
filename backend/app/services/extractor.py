"""Multi-format text extraction.

``extract_blocks`` turns a raw upload into an ordered list of text *blocks*,
each carrying its provenance (page number, section heading, heading level).
Downstream, the chunker concatenates blocks and inherits that provenance, so
every chunk knows the page and section it came from.

Supported types:
  * pdf        — PyMuPDF per page; if a page yields < 50 chars of embedded text
                 it is OCR'd (best-effort; degrades to empty text if the
                 tesseract binary is unavailable).
  * docx       — python-docx paragraphs, preserving heading styles/levels.
  * pptx       — python-pptx: one "page" per slide, slide title as a heading,
                 body text frames as body blocks, speaker notes tagged
                 section_heading="Notes".
  * xlsx       — openpyxl; each row rendered as "Column: Value, ..." text.
  * csv        — same row rendering as xlsx.
  * txt / md   — split on blank lines.

Anything else raises HTTPException(415).
"""
from __future__ import annotations

import csv as _csv
import io
import logging

from fastapi import HTTPException

logger = logging.getLogger("papertrail.extractor")

# Types this service can turn into text.
SUPPORTED_TYPES = {"pdf", "docx", "pptx", "txt", "md", "xlsx", "csv"}

# Below this many characters of embedded text, a PDF page is treated as scanned
# and sent to OCR.
_OCR_MIN_CHARS = 50

# A "block" is a dict with at least "text"; optional provenance keys:
#   page (int), heading (str|None), level (int 0-3), is_heading (bool).
Block = dict


def _block(text: str, *, page: int = 1, heading: str | None = None,
           level: int = 0, is_heading: bool = False) -> Block:
    return {
        "text": text,
        "page": page,
        "heading": heading,
        "level": level,
        "is_heading": is_heading,
    }


# ------------------------------- PDF ----------------------------------- #
def _ocr_pixmap(page) -> str:
    """Best-effort OCR of a rendered PDF page. Returns "" if OCR is unavailable."""
    try:
        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(dpi=200)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        return pytesseract.image_to_string(img).strip()
    except Exception as exc:  # noqa: BLE001 - OCR is optional; never fail the upload
        logger.warning("OCR skipped (tesseract unavailable or failed): %s", exc)
        return ""


def _extract_pdf(data: bytes) -> list[Block]:
    import fitz  # PyMuPDF

    blocks: list[Block] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for i, page in enumerate(doc, start=1):
            text = (page.get_text("text") or "").strip()
            if len(text) < _OCR_MIN_CHARS:
                ocr = _ocr_pixmap(page)
                if len(ocr) > len(text):
                    text = ocr
            if text:
                blocks.append(_block(text, page=i))
    return blocks


# ------------------------------- DOCX ---------------------------------- #
def _docx_heading_level(style_name: str | None) -> int:
    """Map a python-docx paragraph style to a heading level 0 (body) or 1-3."""
    if not style_name:
        return 0
    name = style_name.lower()
    if name.startswith("title"):
        return 1
    if name.startswith("heading"):
        digits = "".join(c for c in name if c.isdigit())
        if digits:
            return min(3, max(1, int(digits)))
        return 2
    return 0


def _extract_docx(data: bytes) -> list[Block]:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    blocks: list[Block] = []
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        level = _docx_heading_level(getattr(para.style, "name", None))
        blocks.append(
            _block(
                text,
                heading=text if level else None,
                level=level,
                is_heading=level > 0,
            )
        )
    return blocks


# ------------------------------- PPTX ----------------------------------- #
def _extract_pptx(data: bytes) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_id = title_shape.shape_id if title_shape is not None else None
        title_text = (title_shape.text or "").strip() if title_shape is not None else ""
        if title_text:
            blocks.append(
                _block(title_text, page=i, heading=title_text, level=1, is_heading=True)
            )

        for shape in slide.shapes:
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if text:
                blocks.append(_block(text, page=i, heading=title_text or None))

        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                blocks.append(_block(notes_text, page=i, heading="Notes"))
    return blocks


# --------------------------- XLSX / CSV -------------------------------- #
def _row_to_text(headers: list[str], values: list) -> str:
    parts = []
    for h, v in zip(headers, values):
        if v is None or v == "":
            continue
        label = str(h).strip() if h not in (None, "") else "Column"
        parts.append(f"{label}: {v}")
    return ", ".join(parts)


def _extract_xlsx(data: bytes) -> list[Block]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[Block] = []
    for sheet in wb.worksheets:
        rows = sheet.iter_rows(values_only=True)
        try:
            header_row = next(rows)
        except StopIteration:
            continue
        headers = [str(c) if c is not None else "" for c in header_row]
        for idx, row in enumerate(rows, start=2):  # row 1 = header
            text = _row_to_text(headers, list(row))
            if text:
                blocks.append(_block(text, heading=sheet.title))
    wb.close()
    return blocks


def _extract_csv(data: bytes) -> list[Block]:
    text = data.decode("utf-8", errors="replace")
    reader = _csv.reader(io.StringIO(text))
    try:
        headers = next(reader)
    except StopIteration:
        return []
    blocks: list[Block] = []
    for row in reader:
        line = _row_to_text(headers, row)
        if line:
            blocks.append(_block(line))
    return blocks


# ------------------------------- TXT ----------------------------------- #
def _extract_txt(data: bytes) -> list[Block]:
    text = data.decode("utf-8", errors="replace")
    parts = [p.strip() for p in text.split("\n\n")]
    return [_block(p) for p in parts if p]


# ------------------------------ facade --------------------------------- #
_EXTRACTORS = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
    "csv": _extract_csv,
    "txt": _extract_txt,
    "md": _extract_txt,
}


def validate_content(data: bytes, file_type: str) -> bool:
    """Best-effort magic-byte / decodability check so a file's real content
    matches its extension (a renamed .exe is rejected before extraction)."""
    ft = file_type.lower()
    head = data[:1024].lstrip()
    if ft == "pdf":
        return head[:5] == b"%PDF-"
    if ft in {"docx", "xlsx", "pptx"}:
        # OOXML files are ZIP archives -> begin with the local-file-header magic.
        return data[:4] in (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")
    if ft in {"txt", "md", "csv"}:
        if b"\x00" in data[:8192]:
            return False
        try:
            data[:8192].decode("utf-8")
            return True
        except UnicodeDecodeError:
            sample = data[:8192]
            printable = sum(
                1 for b in sample if b in (9, 10, 13) or 32 <= b <= 126 or b >= 160
            )
            return printable / max(1, len(sample)) > 0.85
    return False


def extract_blocks(data: bytes, file_type: str) -> list[Block]:
    """Extract ordered text blocks from ``data`` for a supported ``file_type``.

    Raises HTTPException(415) for unsupported types.
    """
    ft = file_type.lower()
    extractor = _EXTRACTORS.get(ft)
    if extractor is None:
        raise HTTPException(status_code=415, detail=f"Unsupported file type '.{ft}'.")
    return extractor(data)


def blocks_to_text(blocks: list[Block]) -> str:
    return "\n\n".join(b["text"] for b in blocks if b.get("text"))


def count_words(text: str) -> int:
    return len(text.split())


def page_count(blocks: list[Block]) -> int:
    """Highest page number seen (>=1 when there is any content)."""
    pages = [b.get("page", 1) for b in blocks]
    return max(pages) if pages else 0
