"""Shared test helpers: tiny in-memory file fixtures and upload utilities."""
from __future__ import annotations

import io


def make_text_pdf(text: str) -> bytes:
    """Build a minimal, valid single-page PDF containing ``text``.

    Hand-assembled with a correct xref table so pypdf reads it without
    recovery — no reportlab dependency required.
    """
    objs = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents 4 0 R"
        b"/Resources<</Font<</F1 5 0 R>>>>>>",
    ]
    stream = b"BT /F1 24 Tf 72 700 Td (" + text.encode() + b") Tj ET"
    objs.append(b"<</Length " + str(len(stream)).encode() + b">>stream\n" + stream + b"\nendstream")
    objs.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")

    out = b"%PDF-1.4\n"
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + o + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n"
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (
        b"trailer<</Size " + str(len(objs) + 1).encode() + b"/Root 1 0 R>>\n"
        b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    )
    return out


def _pdf_escape(s: str) -> bytes:
    return s.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)").encode()


def make_multi_page_pdf(page_texts: list[str]) -> bytes:
    """Build a minimal, valid multi-page PDF with one paragraph per page.

    Generalizes ``make_text_pdf`` to N pages sharing a single font object, with
    long page text word-wrapped across several text-showing lines so pypdf's
    extract_text() returns continuous prose rather than one huge line.
    """
    n = len(page_texts)
    font_obj = 3 + 2 * n

    objs: list[bytes] = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids["
        + b" ".join(f"{3 + 2 * i} 0 R".encode() for i in range(n))
        + b"]/Count "
        + str(n).encode()
        + b">>",
    ]
    for i, text in enumerate(page_texts):
        page_obj = 3 + 2 * i
        content_obj = page_obj + 1
        objs.append(
            b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents "
            + str(content_obj).encode()
            + b" 0 R/Resources<</Font<</F1 "
            + str(font_obj).encode()
            + b" 0 R>>>>>>"
        )

        words = text.split()
        lines: list[str] = []
        cur = ""
        for w in words:
            if cur and len(cur) + len(w) + 1 > 90:
                lines.append(cur)
                cur = w
            else:
                cur = f"{cur} {w}".strip()
        if cur:
            lines.append(cur)

        parts = [b"BT /F1 12 Tf 50 750 Td"]
        for j, line in enumerate(lines):
            if j == 0:
                parts.append(b" (" + _pdf_escape(line) + b") Tj")
            else:
                parts.append(b" 0 -16 Td (" + _pdf_escape(line) + b") Tj")
        parts.append(b" ET")
        stream = b"".join(parts)
        objs.append(
            b"<</Length " + str(len(stream)).encode() + b">>stream\n" + stream + b"\nendstream"
        )
    objs.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")

    out = b"%PDF-1.4\n"
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(out))
        out += str(i).encode() + b" 0 obj\n" + o + b"\nendobj\n"
    xref_pos = len(out)
    out += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n"
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += ("%010d 00000 n \n" % off).encode()
    out += (
        b"trailer<</Size " + str(len(objs) + 1).encode() + b"/Root 1 0 R>>\n"
        b"startxref\n" + str(xref_pos).encode() + b"\n%%EOF"
    )
    return out


def make_pptx_bytes(slides: list[dict]) -> bytes:
    """Build a real .pptx from a list of {"title", "body", "notes"} dicts
    (each key optional) using python-pptx's default "Title and Content" layout."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    for s in slides:
        slide = prs.slides.add_slide(layout)
        if s.get("title"):
            slide.shapes.title.text = s["title"]
        if s.get("body"):
            slide.placeholders[1].text_frame.text = s["body"]
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_encrypted_pdf(password: str = "secret") -> bytes:
    """A password-protected PDF whose text cannot be extracted without the key."""
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.encrypt(password)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def upload_bytes(client, filename: str, data: bytes, content_type: str = "application/octet-stream"):
    """POST a file to the upload endpoint and return the raw response."""
    return client.post(
        "/api/documents/upload",
        files={"file": (filename, data, content_type)},
    )


def upload_text_doc(client, filename: str, text: str):
    """Upload a plain-text document and assert success; return the JSON body."""
    res = upload_bytes(client, filename, text.encode("utf-8"), "text/plain")
    assert res.status_code == 200, res.text
    return res.json()
